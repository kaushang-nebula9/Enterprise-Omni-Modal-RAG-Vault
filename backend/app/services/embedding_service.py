"""
Embedding service using sentence-transformers for text embeddings, 
audio transcription, and slide image description.
"""
import time
import logging
from google import genai
from google.genai import types as genai_types
from app.core.config import settings

logger = logging.getLogger(__name__)

_client: genai.Client | None = None


def _get_client() -> genai.Client:
    """Lazily initialise and return the Gemini client."""
    global _client
    if _client is None:
        _client = genai.Client(api_key=settings.GEMINI_API_KEY)
    return _client


# ---------------------------------------------------------------------------
# Text Embedding
# ---------------------------------------------------------------------------

# -- Sentence Transformers (ACTIVE) -----------------------------------------
# Model: BAAI/bge-large-en-v1.5  |  Output dimension: 1024

from sentence_transformers import SentenceTransformer as _SentenceTransformer

_st_model: _SentenceTransformer | None = None


def _get_st_model() -> _SentenceTransformer:
    """Lazily load and cache the SentenceTransformer model."""
    global _st_model
    if _st_model is None:
        logger.info("Loading SentenceTransformer model: BAAI/bge-large-en-v1.5")
        _st_model = _SentenceTransformer("BAAI/bge-large-en-v1.5")
    return _st_model


def embed_text(text: str) -> list[float]:
    """
    Generate a 1024-dimensional text embedding using BAAI/bge-large-en-v1.5
    via Sentence Transformers (local, no API calls).

    Raises an exception with a clear message if the call fails.
    """
    try:
        model = _get_st_model()
        vector = model.encode(text, normalize_embeddings=True)
        return vector.tolist()
    except Exception as exc:
        raise RuntimeError(f"Failed to generate embedding: {exc}") from exc


# ---------------------------------------------------------------------------

def transcribe_audio(file_path: str) -> str:
    """
    Transcribe an audio file using Gemini 2.5 Flash multimodal.

    Loads the audio from the given absolute path, sends it to the model with a
    transcription prompt, and returns the transcription text.
    Raises an exception with a clear message if transcription fails.
    """
    client = _get_client()
    try:
        with open(file_path, "rb") as f:
            audio_bytes = f.read()

        # Determine MIME type from extension
        ext = file_path.rsplit(".", 1)[-1].lower()
        mime_map = {
            "mp3": "audio/mpeg",
            "wav": "audio/wav",
            "m4a": "audio/mp4",
            "webm": "audio/webm",
            "ogg": "audio/ogg",
        }
        mime_type = mime_map.get(ext, "audio/mpeg")

        response = client.models.generate_content(
            model="gemini-2.5-flash",
            contents=[
                genai_types.Part.from_bytes(data=audio_bytes, mime_type=mime_type),
                "Transcribe this audio file accurately. Return only the transcription text, nothing else.",
            ],
        )
        return response.text or ""
    except Exception as exc:
        raise RuntimeError(f"Failed to transcribe audio '{file_path}': {exc}") from exc


def process_pptx_slides(file_path: str) -> list[dict]:
    """
    Extract text, tables, charts, and images from a PPTX file using python-pptx.
    Images are passed to Gemini Vision for description.
    Merges content per slide, splits if too long, and generates embeddings.
    
    Returns a list of dicts: [{"slide_number": int, "slide_title": str, "text": str, "vector": list[float]}]
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    

    try:
        prs = Presentation(file_path)
    except Exception as exc:
        logger.warning("Failed to open PPTX file '%s': %s", file_path, exc)
        return []

    results = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_title = ""
        if slide.shapes.title and slide.shapes.title.text:
            slide_title = slide.shapes.title.text.strip()

        texts = []
        for shape in slide.shapes:
            # Text frames (titles, bullet points)
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        texts.append(line)
            # Tables
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        texts.append(row_text)
            
            # Charts
            if shape.has_chart:
                try:
                    chart = shape.chart
                    c_title = chart.chart_title.text_frame.text if chart.has_title else "Untitled Chart"
                    c_type = str(chart.chart_type).split('.')[-1]
                    texts.append(f"Chart: {c_title}, Type: {c_type}")
                    
                    for series in chart.series:
                        s_name = getattr(series, "name", "Series")
                        cats = []
                        try:
                            cats = [c.label for c in chart.plots[0].categories]
                        except Exception:
                            pass
                        
                        vals = series.values
                        if cats and len(cats) == len(vals):
                            data_str = ", ".join(f"{c}: {v}" for c, v in zip(cats, vals))
                        else:
                            data_str = ", ".join(str(v) for v in vals)
                        
                        texts.append(f"Series [{s_name}]: {data_str}")
                except Exception as exc:
                    logger.warning("Failed to extract chart on slide %d: %s", slide_idx, exc)

            # Images
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image_bytes = shape.image.blob
                    ext = shape.image.ext
                    import base64
                    from anthropic import Anthropic
                    from app.core.config import settings

                    anthropic_client = Anthropic(api_key=settings.ANTHROPIC_API_KEY)
                    
                    mime_map = {
                        "png": "image/png",
                        "jpg": "image/jpeg",
                        "jpeg": "image/jpeg",
                        "gif": "image/gif",
                        "webp": "image/webp",
                    }
                    ext_lower = (ext or "").lower()
                    mime_type = mime_map.get(ext_lower, "image/png")
                    
                    b64_image = base64.b64encode(image_bytes).decode("utf-8")
                    system_prompt = "If the image is a chart or a graph, then give exact values and analyze the data in a structured format."
                    
                    message = anthropic_client.messages.create(
                        model="claude-haiku-4-5-20251001",
                        max_tokens=4096,
                        temperature=0,
                        system=system_prompt,
                        messages=[
                            {
                                "role": "user",
                                "content": [
                                    {
                                        "type": "image",
                                        "source": {
                                            "type": "base64",
                                            "media_type": mime_type,
                                            "data": b64_image,
                                        },
                                    },
                                    {
                                        "type": "text",
                                        "text": "Describe this image in detail. Be specific about data, trends, and key takeaways."
                                    }
                                ],
                            }
                        ]
                    )
                    desc = message.content[0].text.strip()
                    if desc:
                        texts.append(f"[Image Description]: {desc}")
                except Exception as exc:
                    logger.warning("Failed to describe image on slide %d: %s", slide_idx, exc)

        # Speaker notes
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                texts.append(f"[Speaker Notes] {notes_text}")

        merged = "\n".join(texts).strip()
        if not merged:
            continue

        # Chunk splitting logic (>2000 tokens approx)
        approx_tokens = len(merged.split()) * 1.3
        chunks = []
        if approx_tokens > 2000 and len(merged) > 1:
            mid = len(merged) // 2
            chunks.append(merged[:mid])
            chunks.append(merged[mid:])
        else:
            chunks.append(merged)
            
        for chunk in chunks:
            try:
                vector = embed_text(chunk)
                results.append({
                    "slide_number": slide_idx,
                    "slide_title": slide_title,
                    "text": chunk,
                    "vector": vector
                })
            except Exception as exc:
                logger.warning("Failed to embed chunk for slide %d: %s", slide_idx, exc)
                
    return results


_anthropic_client = None


def _get_anthropic_client():
    """Lazily initialise and return the Anthropic client."""
    global _anthropic_client
    if _anthropic_client is None:
        from anthropic import Anthropic
        _anthropic_client = Anthropic(api_key=settings.ANTHROPIC_API_KEY)
    return _anthropic_client


def generate_document_description(content_sample: str, file_type: str) -> str | None:
    """
    Generate a single concise sentence (max 50 words) describing the document's content.
    Uses Claude to generate the description best-effort.
    """
    try:
        client = _get_anthropic_client()
        prompt = (
            "Based on the following content from a document, write a single concise sentence "
            "(max 50 words) describing what this document is about. Do not mention the file format. "
            "Be specific about the topic or purpose.\n\n"
            f"Content: {content_sample}\n\n"
            "Description:"
        )
        message = client.messages.create(
            model="claude-haiku-4-5-20251001",
            max_tokens=150,
            temperature=0,
            messages=[
                {"role": "user", "content": prompt}
            ]
        )
        return (message.content[0].text or "").strip()
    except Exception as exc:
        logger.error("Failed to generate document description: %s", exc)
        return None

